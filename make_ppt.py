from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ────────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x1B, 0x2A)   # 네이비
BG_CARD    = RGBColor(0x1B, 0x2B, 0x3E)   # 카드 배경
ACCENT     = RGBColor(0x4F, 0xC3, 0xF7)   # 하늘색 강조
ACCENT2    = RGBColor(0x81, 0xE6, 0xD9)   # 민트
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xB0, 0xBE, 0xC5)
YELLOW     = RGBColor(0xFF, 0xD7, 0x54)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # 상단 accent 바
    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

    # 메인 타이틀
    add_text(slide, "AI 논문 데일리 리뷰", Inches(1), Inches(1.5), Inches(11), Inches(1.4),
             size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "자동화 시스템", Inches(1), Inches(2.8), Inches(11), Inches(1.2),
             size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # 부제
    add_text(slide, "매일 아침, AI 최신 논문 3편이 내 메일함에 도착한다",
             Inches(1), Inches(4.2), Inches(11), Inches(0.8),
             size=22, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, "kimheekyoung  |  2026",
             Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), ACCENT)
    return slide


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

    add_text(slide, "😓  문제", Inches(0.6), Inches(0.3), Inches(6), Inches(0.6),
             size=18, color=ACCENT)
    add_text(slide, "논문이 너무 많고, 너무 어렵다",
             Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
             size=38, bold=True, color=WHITE)

    items = [
        ("📄", "매일 수백 편의 AI 논문이 쏟아진다"),
        ("⏱️", "논문 1편 읽는 데 평균 2~3시간"),
        ("🔤", "영어 전문 용어 → 이해하기 어렵다"),
        ("😴", "바쁜 일상 속 논문 읽을 시간이 없다"),
    ]
    for i, (icon, text) in enumerate(items):
        y = Inches(2.1) + i * Inches(1.1)
        add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.9), BG_CARD)
        add_text(slide, icon, Inches(0.9), y + Pt(8), Inches(0.6), Inches(0.8), size=24)
        add_text(slide, text, Inches(1.7), y + Pt(8), Inches(10.5), Inches(0.8), size=22, color=WHITE)

    return slide


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT2)

    add_text(slide, "💡  해결책", Inches(0.6), Inches(0.3), Inches(6), Inches(0.6),
             size=18, color=ACCENT2)
    add_text(slide, "매일 오전 9시, 자동으로 배달되는 논문 리뷰",
             Inches(0.6), Inches(0.9), Inches(12), Inches(0.9),
             size=32, bold=True, color=WHITE)

    steps = [
        ("①", "오늘 가장 핫한 AI 논문 3편 자동 수집"),
        ("②", "논문 전문을 AI가 직접 읽는다"),
        ("③", "고등학생도 이해할 수 있게 쉽게 정리"),
        ("④", "아침에 일어나면 메일함에 도착!"),
    ]
    for i, (num, text) in enumerate(steps):
        y = Inches(2.1) + i * Inches(1.05)
        add_rect(slide, Inches(0.6), y, Inches(0.75), Inches(0.85), ACCENT2)
        add_text(slide, num, Inches(0.6), y, Inches(0.75), Inches(0.85),
                 size=22, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, text, Inches(1.5), y + Pt(6), Inches(10.8), Inches(0.75),
                 size=22, color=WHITE)

    return slide


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

    add_text(slide, "전체 파이프라인",
             Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    boxes = [
        (ACCENT,  "🌐 HuggingFace",  "트렌딩 논문 수집"),
        (ACCENT2, "📖 Jina Reader",   "논문 전문 fetch"),
        (YELLOW,  "🤖 GitHub Models", "AI 리뷰 생성"),
        (RGBColor(0xFC, 0x8B, 0x5E), "📧 Gmail", "이메일 발송"),
    ]

    bw, bh = Inches(2.6), Inches(2.2)
    total = len(boxes)
    gap = (W - bw * total) / (total + 1)

    for i, (color, title, sub) in enumerate(boxes):
        x = gap + i * (bw + gap)
        y = Inches(1.5)
        add_rect(slide, x, y, bw, bh, color)
        add_text(slide, title, x, y + Inches(0.3), bw, Inches(0.8),
                 size=20, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, sub, x, y + Inches(1.0), bw, Inches(0.7),
                 size=16, color=BG_DARK, align=PP_ALIGN.CENTER)

        if i < total - 1:
            ax = x + bw + gap * 0.2
            ay = y + bh / 2 - Inches(0.15)
            add_text(slide, "→", ax, ay, gap * 0.6, Inches(0.4),
                     size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.08), GRAY)
    add_text(slide, "⏰  매일 00:00 UTC (오전 9시 KST)  GitHub Actions가 자동 실행",
             Inches(0.6), Inches(4.3), Inches(12.1), Inches(0.6),
             size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(slide, "한 번 설정하면 → 매일 자동으로 동작!",
             Inches(0.6), Inches(5.1), Inches(12.1), Inches(0.6),
             size=22, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    return slide


def slide_step(prs, step_num, title, icon, color, points, analogy):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), color)

    add_text(slide, f"STEP {step_num}", Inches(0.6), Inches(0.3), Inches(3), Inches(0.5),
             size=16, color=color)
    add_text(slide, f"{icon}  {title}",
             Inches(0.6), Inches(0.8), Inches(12), Inches(0.8),
             size=34, bold=True, color=WHITE)

    for i, (label, desc) in enumerate(points):
        y = Inches(1.9) + i * Inches(1.0)
        add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.85), BG_CARD)
        add_text(slide, label, Inches(0.8), y + Pt(5), Inches(1.8), Inches(0.75),
                 size=15, bold=True, color=color)
        add_text(slide, desc, Inches(2.5), y + Pt(5), Inches(9.8), Inches(0.75),
                 size=19, color=WHITE)

    add_rect(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.3), RGBColor(0x1A, 0x35, 0x50))
    add_text(slide, "💡 쉽게 말하면", Inches(0.9), Inches(5.55), Inches(4), Inches(0.4),
             size=14, color=ACCENT)
    add_text(slide, analogy, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.7),
             size=19, color=WHITE, bold=True)

    return slide


def slide_github_actions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), YELLOW)

    add_text(slide, "⚙️  자동화", Inches(0.6), Inches(0.3), Inches(4), Inches(0.5),
             size=16, color=YELLOW)
    add_text(slide, "GitHub Actions — 내 컴퓨터가 꺼져 있어도 동작!",
             Inches(0.6), Inches(0.8), Inches(12), Inches(0.8),
             size=30, bold=True, color=WHITE)

    add_rect(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), BG_CARD)
    add_text(slide, "설정 파일 (daily_review.yml)",
             Inches(0.8), Inches(1.9), Inches(5.4), Inches(0.5),
             size=16, bold=True, color=ACCENT)
    code = """on:
  schedule:
    - cron: '0 0 * * *'
  workflow_dispatch:

jobs:
  send-daily-review:
    permissions:
      models: read
    steps:
      - run: python daily_review.py"""
    add_text(slide, code, Inches(0.8), Inches(2.4), Inches(5.4), Inches(3.8),
             size=14, color=ACCENT2)

    secrets = [
        ("✅", "GMAIL_USER",         "발신 Gmail 주소"),
        ("✅", "GMAIL_APP_PASSWORD", "Gmail 앱 비밀번호"),
        ("🎁", "GITHUB_TOKEN",       "자동 제공! 등록 불필요"),
    ]
    add_text(slide, "필요한 설정", Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.5),
             size=16, bold=True, color=YELLOW)
    for i, (icon, key, desc) in enumerate(secrets):
        y = Inches(2.5) + i * Inches(1.0)
        add_rect(slide, Inches(7.0), y, Inches(5.7), Inches(0.85), BG_CARD)
        add_text(slide, f"{icon} {key}", Inches(7.2), y + Pt(5), Inches(3.5), Inches(0.4),
                 size=17, bold=True, color=YELLOW)
        add_text(slide, desc, Inches(7.2), y + Pt(30), Inches(5.3), Inches(0.4),
                 size=14, color=GRAY)

    return slide


def slide_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

    add_text(slide, "기술 스택 & 비용",
             Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("논문 수집",   "HuggingFace Daily Papers API", "무료"),
        ("전문 fetch",  "Jina Reader (r.jina.ai)",      "무료"),
        ("AI 리뷰 생성","GitHub Models (gpt-4o-mini)",   "무료"),
        ("이메일 발송", "Gmail SMTP SSL",                "무료"),
        ("자동화",      "GitHub Actions",                "무료"),
    ]
    headers = ["역할", "기술", "비용"]
    col_w = [Inches(2.8), Inches(6.5), Inches(2.3)]
    col_x = [Inches(0.4), Inches(3.2), Inches(9.7)]

    # 헤더
    for j, (hdr, cw, cx) in enumerate(zip(headers, col_w, col_x)):
        add_rect(slide, cx, Inches(1.3), cw - Inches(0.1), Inches(0.55), ACCENT)
        add_text(slide, hdr, cx + Inches(0.05), Inches(1.32), cw - Inches(0.15), Inches(0.5),
                 size=17, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        y = Inches(1.85) + i * Inches(0.95)
        bg = BG_CARD if i % 2 == 0 else BG_DARK
        add_rect(slide, Inches(0.4), y, Inches(12.6), Inches(0.88), bg)
        for j, (val, cw, cx) in enumerate(zip(row, col_w, col_x)):
            color = YELLOW if j == 2 else WHITE
            add_text(slide, val, cx + Inches(0.1), y + Pt(8), cw - Inches(0.15), Inches(0.7),
                     size=18, color=color, align=PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT)

    add_text(slide, "💰  총 비용: 0원  (GitHub 계정 + Gmail만 있으면 OK!)",
             Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.5),
             size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

    return slide


def slide_ending(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

    add_text(slide, "이 프로젝트가 해결하는 것",
             Inches(0.6), Inches(0.5), Inches(12), Inches(0.7),
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    checkmarks = [
        "매일 최신 AI 트렌드를 5분 안에 파악",
        "어려운 논문을 쉬운 언어로 이해",
        "완전 자동화 — 한 번 설정하면 끝",
        "비용 0원",
    ]
    for i, txt in enumerate(checkmarks):
        y = Inches(1.5) + i * Inches(1.0)
        add_rect(slide, Inches(1.5), y, Inches(10.3), Inches(0.82), BG_CARD)
        add_text(slide, "✅", Inches(1.7), y + Pt(6), Inches(0.7), Inches(0.7), size=22)
        add_text(slide, txt, Inches(2.5), y + Pt(6), Inches(9.0), Inches(0.7),
                 size=22, color=WHITE)

    add_rect(slide, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.1),
             RGBColor(0x0A, 0x28, 0x47))
    add_text(slide, '"AI가 AI 논문을 읽고, 사람이 이해할 수 있게 설명해준다"',
             Inches(0.8), Inches(5.75), Inches(11.7), Inches(0.9),
             size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(slide, "GitHub: heezzing/Today_research_paper",
             Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35),
             size=14, color=GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), ACCENT)
    return slide


# ── 슬라이드 생성 ──────────────────────────────────────────────────────────────

prs = new_prs()

slide_title(prs)
slide_problem(prs)
slide_solution(prs)
slide_pipeline(prs)

slide_step(prs, 1, "트렌딩 논문 수집", "🌐", ACCENT,
    [
        ("출처",  "HuggingFace Daily Papers (huggingface.co/papers)"),
        ("방법",  "전 세계 AI 연구자들이 '좋아요'를 누른 논문 순위"),
        ("결과",  "오늘 가장 주목받는 논문 상위 3편 자동 선정"),
    ],
    "마치 유튜브 인기 급상승 영상처럼, 오늘 AI 커뮤니티가 가장 주목한 논문을 뽑아준다")

slide_step(prs, 2, "논문 전문 가져오기", "📖", ACCENT2,
    [
        ("도구",  "Jina Reader — r.jina.ai"),
        ("방법",  "arxiv PDF 링크를 넣으면 읽기 쉬운 텍스트로 변환"),
        ("결과",  "AI가 읽을 수 있는 형태로 논문 전문 준비 완료"),
    ],
    "마치 번역 앱이 사진 속 글자를 읽어주듯, PDF 논문을 텍스트로 바꿔준다")

slide_step(prs, 3, "AI가 리뷰를 써준다", "🤖", YELLOW,
    [
        ("모델",  "GitHub Models — gpt-4o-mini (완전 무료!)"),
        ("입력",  "논문 제목 + 저자 + 초록 + 전문"),
        ("출력",  "6개 섹션으로 구성된 한국어 리뷰 문서"),
    ],
    "마치 똑똑한 친구가 어려운 책을 읽고 쉽게 요약해주는 것처럼")

slide_step(prs, 4, "이메일 자동 발송", "📧", RGBColor(0xFC, 0x8B, 0x5E),
    [
        ("방법",  "Gmail SMTP SSL (포트 465) — 안전한 암호화 전송"),
        ("형식",  "논문 3편 리뷰를 하나의 이메일로 합산"),
        ("시각",  "매일 오전 9시 (KST) 자동 발송"),
    ],
    "설정 한 번이면 매일 아침 자동으로 신문처럼 배달된다")

slide_github_actions(prs)
slide_stack(prs)
slide_ending(prs)

out = "/Users/hk/study/mcp-servers/AI_논문_데일리리뷰_발표.pptx"
prs.save(out)
print(f"저장 완료: {out}")
